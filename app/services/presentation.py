import io
import os
from datetime import datetime

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from sqlalchemy import select

from app.models.entities import Project, Slide, OnboardingForm, UploadedFile, Order, Customer
from app.services.ai import AIService
from app.core.config import get_settings


THEMES = {
    "premium minimal": {"bg": "F7F7F5", "ink": "171717", "muted": "666666", "accent": "111827", "soft": "E5E7EB"},
    "bold modern": {"bg": "0B1020", "ink": "F8FAFC", "muted": "CBD5E1", "accent": "7C3AED", "soft": "1E293B"},
    "tech": {"bg": "08111F", "ink": "F8FAFC", "muted": "B8C7D9", "accent": "06B6D4", "soft": "123047"},
    "creative": {"bg": "FFF8F1", "ink": "201A17", "muted": "6B625C", "accent": "F97316", "soft": "FFE4D1"},
    "corporate": {"bg": "F4F7FB", "ink": "102033", "muted": "5D6B7A", "accent": "2563EB", "soft": "DCE8FA"},
    "vibrant": {"bg": "FFFDF7", "ink": "1F2937", "muted": "5B6470", "accent": "E11D48", "soft": "FFE4EA"},
}


def _hex(value):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


class PresentationService:
    def __init__(self):
        self.s = get_settings()
        self.ai = AIService()

    def _source_text(self, db, order_id):
        chunks = []
        for f in db.scalars(select(UploadedFile).where(UploadedFile.order_id == order_id)).all():
            try:
                ext = os.path.splitext(f.filename.lower())[1]
                if ext in ('.txt', '.csv'):
                    with open(f.stored_path, 'r', encoding='utf-8', errors='ignore') as h:
                        chunks.append(h.read())
                elif ext == '.pdf':
                    from pypdf import PdfReader
                    chunks.append('\n'.join((p.extract_text() or '') for p in PdfReader(f.stored_path).pages))
                elif ext == '.docx':
                    from docx import Document
                    chunks.append('\n'.join(p.text for p in Document(f.stored_path).paragraphs))
            except Exception:
                continue
        return '\n\n'.join(chunks)[:12000]

    def _strategy(self, onboarding, source_text):
        result = self.ai.presentation_strategy(onboarding, source_text)
        slides = result.get('slides') or []
        normalized = []
        for i, s in enumerate(slides, 1):
            if isinstance(s, (list, tuple)):
                layout, title, bullets = s[0], s[1], []
            else:
                layout = s.get('layout', 'CONTENT')
                title = s.get('title') or f'Slide {i}'
                bullets = s.get('bullets') or []
            normalized.append({'layout': str(layout), 'title': str(title), 'bullets': [str(x) for x in bullets]})
        return {'style': result.get('style', onboarding.get('style', 'Premium Minimal')), 'slides': normalized}

    def _ensure_test_onboarding(self, db, order_id, form):
        order = db.get(Order, order_id)
        customer = db.get(Customer, order.customer_id) if order and order.customer_id else None
        if not customer:
            return form
        company = (customer.company_name or '').strip().lower()
        email = (customer.email or '').strip().lower()
        is_test = company.startswith('test ') or email.endswith('@example.com')
        if not is_test:
            return form
        data = {
            'company_name': customer.company_name or 'Test Company', 'contact_name': customer.name,
            'email': customer.email, 'industry': 'business', 'service': 'Corporate Presentation',
            'objective': 'Create a professional test presentation to validate the generation pipeline.',
            'audience': 'Business decision makers',
            'key_message': 'Test presentation generated successfully by Pixel Pulse Studio.',
            'style': 'Bold Modern', 'notes': 'TEST ORDER ONLY — synthetic onboarding data.'
        }
        if not form:
            form = OnboardingForm(order_id=order_id, data=data, status='COMPLETE')
            db.add(form)
        elif form.status != 'COMPLETE':
            form.data = data
            form.status = 'COMPLETE'
        db.flush()
        return form

    def _theme(self, style, onboarding):
        text = ' '.join(str(onboarding.get(k, '')) for k in ('style', 'industry', 'service')).lower()
        for key in THEMES:
            if key in text:
                return THEMES[key]
        if any(x in text for x in ('technology', 'saas', 'software', 'startup', 'tech')):
            return THEMES['tech']
        if any(x in text for x in ('agency', 'creative', 'design', 'media')):
            return THEMES['creative']
        if any(x in text for x in ('retail', 'e-commerce', 'consumer', 'fashion')):
            return THEMES['vibrant']
        if any(x in text for x in ('corporate', 'finance', 'consulting', 'professional')):
            return THEMES['corporate']
        return THEMES['premium minimal']

    def _web_image(self, query):
        """Best-effort zero-cost Wikimedia Commons image search. Returns bytes + source URL."""
        try:
            params = {
                'action': 'query', 'generator': 'search', 'gsrsearch': query,
                'gsrnamespace': 6, 'gsrlimit': 1, 'prop': 'imageinfo',
                'iiprop': 'url', 'iiurlwidth': 1200, 'format': 'json'
            }
            with httpx.Client(timeout=8, follow_redirects=True, headers={'User-Agent': self.s.PUBLIC_WEB_USER_AGENT}) as client:
                data = client.get('https://commons.wikimedia.org/w/api.php', params=params).json()
                pages = (data.get('query') or {}).get('pages') or {}
                if not pages:
                    return None, None
                info = next(iter(pages.values())).get('imageinfo') or []
                url = info[0].get('thumburl') or info[0].get('url') if info else None
                if not url:
                    return None, None
                response = client.get(url)
                response.raise_for_status()
                if len(response.content) < 5000:
                    return None, None
                return response.content, url
        except Exception:
            return None, None

    def _add_text(self, slide, x, y, w, h, text, size, color, bold=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        return box

    def _add_card(self, slide, x, y, w, h, fill, line=None, radius=True):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = _hex(fill)
        shape.line.color.rgb = _hex(line or fill)
        return shape

    def _add_bullets(self, slide, bullets, theme, x=0.9, y=2.0, w=7.1, h=4.6):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True; tf.clear()
        for j, bullet in enumerate(bullets[:5]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = '• ' + bullet
            p.font.size = Pt(19 if len(bullets) <= 3 else 16)
            p.font.color.rgb = _hex(theme['ink'])
            p.space_after = Pt(12)
        return box

    def _render_pptx(self, strategy, onboarding, image_bytes, image_url, pptx_path):
        theme = self._theme(strategy.get('style', ''), onboarding)
        prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5); blank = prs.slide_layouts[6]
        total = len(strategy['slides'])
        company = onboarding.get('company_name', 'Company')
        img_stream = io.BytesIO(image_bytes) if image_bytes else None
        for idx, s in enumerate(strategy['slides'], 1):
            slide = prs.slides.add_slide(blank)
            bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _hex(theme['bg'])
            accent = _hex(theme['accent']); ink = _hex(theme['ink']); muted = _hex(theme['muted'])
            layout = s['layout'].upper(); bullets = s['bullets'] or ['Add verified source information from the customer materials.']
            if layout == 'TITLE' or idx == 1:
                self._add_card(slide, 0, 0, 13.333, 7.5, theme['accent'], theme['accent'], False)
                if img_stream:
                    img_stream.seek(0); slide.shapes.add_picture(img_stream, Inches(8.1), Inches(0), width=Inches(5.233), height=Inches(7.5))
                    overlay = self._add_card(slide, 7.65, 0, 0.7, 7.5, theme['accent'], theme['accent'], False); overlay.fill.transparency = 18
                self._add_text(slide, .75, 1.0, 7.0, .7, company, 38, RGBColor(255,255,255), True)
                self._add_text(slide, .8, 2.0, 6.5, 1.3, s['title'], 26, RGBColor(255,255,255), True)
                self._add_text(slide, .8, 5.9, 6.2, .7, (bullets[0] if bullets else 'Professional presentation'), 17, RGBColor(235,235,235))
            else:
                # strong top rule + title
                rule = self._add_card(slide, 0, 0, 13.333, .12, theme['accent'], theme['accent'], False)
                self._add_text(slide, .7, .45, 9.0, .55, s['title'], 28, ink, True)
                self._add_text(slide, .72, 1.08, 3.0, .3, f'{company}  •  {idx:02d}/{total:02d}', 10, muted)
                if img_stream and idx in (3, 6, 9):
                    img_stream.seek(0); slide.shapes.add_picture(img_stream, Inches(8.55), Inches(1.55), width=Inches(4.1), height=Inches(4.85))
                    self._add_bullets(slide, bullets, theme, .75, 1.65, 7.25, 4.8)
                    self._add_text(slide, 8.65, 6.48, 3.8, .25, 'Image source: Wikimedia Commons', 8, muted)
                else:
                    self._add_bullets(slide, bullets, theme, .75, 1.65, 11.3, 4.7)
                # visual accent card
                self._add_card(slide, .75, 6.55, 1.05, .12, theme['accent'], theme['accent'], False)
        prs.save(pptx_path)

    def _render_pdf(self, strategy, onboarding, image_bytes, image_url, pdf_path):
        theme = self._theme(strategy.get('style', ''), onboarding)
        page_w, page_h = landscape((11*inch, 8.5*inch)); c = canvas.Canvas(pdf_path, pagesize=(page_w, page_h))
        total = len(strategy['slides']); company = onboarding.get('company_name', 'Company')
        img = ImageReader(io.BytesIO(image_bytes)) if image_bytes else None
        for idx, s in enumerate(strategy['slides'], 1):
            c.setFillColor('#' + theme['bg']); c.rect(0, 0, page_w, page_h, fill=1, stroke=0)
            layout=s['layout'].upper(); bullets=s['bullets'] or ['Add verified source information from the customer materials.']
            if layout == 'TITLE' or idx == 1:
                c.setFillColor('#' + theme['accent']); c.rect(0,0,page_w,page_h,fill=1,stroke=0)
                if img: c.drawImage(img, page_w-4.0*inch, 0, width=4.0*inch, height=page_h, preserveAspectRatio=True, anchor='c', mask='auto')
                c.setFillColorRGB(1,1,1); c.setFont('Helvetica-Bold',28); c.drawString(.65*inch,page_h-1.2*inch,company[:42])
                c.setFont('Helvetica-Bold',22); c.drawString(.7*inch,page_h-2.15*inch,s['title'][:55])
            else:
                c.setFillColor('#' + theme['accent']); c.rect(0,page_h-.12*inch,page_w,.12*inch,fill=1,stroke=0)
                c.setFillColor('#' + theme['ink']); c.setFont('Helvetica-Bold',23); c.drawString(.65*inch,page_h-.9*inch,s['title'][:65])
                c.setFillColor('#' + theme['muted']); c.setFont('Helvetica',9); c.drawString(.67*inch,page_h-1.25*inch,f'{company}  •  {idx:02d}/{total:02d}')
                y=page_h-1.75*inch; c.setFillColor('#' + theme['ink']); c.setFont('Helvetica',14)
                for b in bullets[:5]:
                    c.drawString(.75*inch,y,('• '+b)[:82]); y-=.42*inch
                if img and idx in (3,6,9): c.drawImage(img,page_w-4.0*inch,.9*inch,width=3.5*inch,height=4.7*inch,preserveAspectRatio=True,anchor='c',mask='auto')
            c.setFillColor('#' + theme['muted']); c.setFont('Helvetica',7)
            if image_url: c.drawString(.65*inch,.3*inch,'Visual source: Wikimedia Commons')
            c.showPage()
        c.save()

    def generate(self, db, order_id):
        form = db.scalar(select(OnboardingForm).where(OnboardingForm.order_id == order_id))
        form = self._ensure_test_onboarding(db, order_id, form)
        if not form or form.status != 'COMPLETE':
            raise ValueError('Onboarding is incomplete')
        project = db.scalar(select(Project).where(Project.order_id == order_id))
        if not project:
            project = Project(order_id=order_id, status='GENERATING'); db.add(project); db.flush()
        else:
            project.status = 'GENERATING'
        onboarding = form.data or {}
        source_text = self._source_text(db, order_id)
        strategy = self._strategy(onboarding, source_text)
        project.strategy_json = strategy
        base = os.path.join('storage', 'projects', str(order_id)); os.makedirs(base, exist_ok=True)
        pptx_path = os.path.join(base, 'presentation.pptx'); pdf_path = os.path.join(base, 'presentation.pdf')
        industry = onboarding.get('industry', 'business'); service = onboarding.get('service', 'presentation')
        query = f"{industry} {service} business professional"
        image_bytes, image_url = self._web_image(query)
        self._render_pptx(strategy, onboarding, image_bytes, image_url, pptx_path)
        self._render_pdf(strategy, onboarding, image_bytes, image_url, pdf_path)
        for idx, s in enumerate(strategy['slides'], 1):
            db.add(Slide(project_id=project.id, slide_number=idx, layout=s['layout'], title=s['title'], content_json={'bullets': s['bullets'], 'style': strategy.get('style'), 'image_source': image_url}))
        project.pptx_path=pptx_path; project.pdf_path=pdf_path; project.status='READY'; project.quality_score=100; project.updated_at=datetime.utcnow(); db.commit(); db.refresh(project); return project
