import os, re, uuid, hashlib
from pathlib import Path
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape, A4
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from app.core.config import get_settings
from app.services.ai import AIService
from app.core.logging import events

LAYOUTS=['TITLE','AGENDA','PROBLEM','SOLUTION','PRODUCT','FEATURES','BENEFITS','MARKET','BUSINESS MODEL','TRACTION','COMPETITION','TEAM','TIMELINE','CASE STUDY','TESTIMONIAL','PRICING','CALL TO ACTION','CONTACT']
class PresentationService:
    def __init__(self): self.s=get_settings(); self.ai=AIService(); self.root=Path('storage'); self.root.mkdir(exist_ok=True)
    def extract_text(self,files):
        chunks=[]
        for p in files:
            ext=Path(p).suffix.lower()
            try:
                if ext=='.txt': chunks.append(Path(p).read_text(errors='ignore'))
                elif ext=='.docx':
                    from docx import Document
                    chunks.append('\n'.join(x.text for x in Document(p).paragraphs))
                elif ext=='.pdf':
                    from pypdf import PdfReader
                    chunks.append('\n'.join(page.extract_text() or '' for page in PdfReader(p).pages))
                elif ext=='.csv': chunks.append(Path(p).read_text(errors='ignore'))
            except Exception as e: events.event('ERROR',component='file_extract',file=str(p),error=str(e))
        return '\n\n'.join(chunks)
    def build(self,order_id,onboarding,files):
        text=self.extract_text(files); strategy=self.ai.presentation_strategy(onboarding,text); slides=strategy.get('slides') or []
        if not slides: slides=[{'layout':'TITLE','title':onboarding.get('presentation_type','Presentation'),'bullets':[]}]
        out=Path('storage/deliveries')/str(order_id); out.mkdir(parents=True,exist_ok=True)
        pptx=out/'presentation.pptx'; pdf=out/'presentation.pdf'
        prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
        for idx,s in enumerate(slides,1): self._slide(prs,s,idx,strategy.get('style','Premium Minimal'))
        prs.save(pptx)
        self._pdf(pdf,slides,strategy.get('style','Premium Minimal'))
        events.event('PRESENTATION_GENERATED',order_id=order_id,slides=len(slides)); return str(pptx),str(pdf),strategy
    def _slide(self,prs,s,idx,style):
        slide=prs.slides.add_slide(prs.slide_layouts[6]); title=str(s.get('title') or f'Slide {idx}'); layout=str(s.get('layout') or 'CONTENT'); bullets=s.get('bullets') or s.get('content') or []
        shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.16)); shape.fill.solid(); shape.line.fill.background()
        tb=slide.shapes.add_textbox(Inches(.8),Inches(.65),Inches(11.8),Inches(1.0)); p=tb.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(28 if idx>1 else 34); p.font.bold=True; p.alignment=PP_ALIGN.LEFT
        body=slide.shapes.add_textbox(Inches(.9),Inches(1.8),Inches(11.3),Inches(4.8)); tf=body.text_frame; tf.word_wrap=True
        if isinstance(bullets,str): bullets=[bullets]
        for i,b in enumerate(bullets[:8]):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=str(b); p.font.size=Pt(18); p.space_after=Pt(10); p.level=0
        foot=slide.shapes.add_textbox(Inches(.8),Inches(7.0),Inches(11.7),Inches(.25)); foot.text_frame.paragraphs[0].text=f'{style}  •  {idx:02d}'; foot.text_frame.paragraphs[0].font.size=Pt(9)
    def _pdf(self,path,slides,style):
        c=canvas.Canvas(str(path),pagesize=landscape(A4)); w,h=landscape(A4)
        for idx,s in enumerate(slides,1):
            c.setFont('Helvetica-Bold',24); c.drawString(45,h-60,str(s.get('title') or f'Slide {idx}'))
            c.setFont('Helvetica',14); y=h-100; bullets=s.get('bullets') or []
            if isinstance(bullets,str): bullets=[bullets]
            for b in bullets[:10]: c.drawString(60,y,str(b)[:110]); y-=24
            c.setFont('Helvetica',8); c.drawString(45,25,f'{style} • {idx:02d}'); c.showPage()
        c.save()
