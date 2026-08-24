from pathlib import Path
from pptx import Presentation
from app.core.config import get_settings
class QCService:
    def __init__(self): self.s=get_settings()
    def inspect(self,pptx,expected_min=1):
        prs=Presentation(pptx); score=100; issues=[]
        if len(prs.slides)<expected_min: score-=25; issues.append('unexpected slide count')
        for i,slide in enumerate(prs.slides,1):
            if not any(getattr(sh,'has_text_frame',False) and sh.text.strip() for sh in slide.shapes): score-=5; issues.append(f'empty slide {i}')
            for sh in slide.shapes:
                if getattr(sh,'has_text_frame',False) and len(sh.text)>1200: score-=3; issues.append(f'overloaded text slide {i}')
        return max(0,score),issues
